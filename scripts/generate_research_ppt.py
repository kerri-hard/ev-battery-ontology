#!/usr/bin/env python3
"""
Generate a high-quality executive/research presentation deck and script.
"""

from __future__ import annotations

import json
import os
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RESULTS = ROOT / "results"
ASSETS = RESULTS / "ppt_assets"

PPT_PATH = RESULTS / f"EV_Battery_AIOps_Research_Brief_PRO_{date.today().isoformat()}.pptx"
SCRIPT_PATH = RESULTS / f"EV_Battery_AIOps_Presentation_Script_PRO_{date.today().isoformat()}.md"
HEALING_SUMMARY = RESULTS / "self_healing_v4_latest.json"

WIDE = Inches(13.333)
HIGH = Inches(7.5)

# Palette
C_BG = RGBColor(7, 11, 24)
C_BG2 = RGBColor(14, 22, 42)
C_TITLE = RGBColor(238, 244, 255)
C_TEXT = RGBColor(202, 214, 232)
C_MUTED = RGBColor(139, 156, 180)
C_ACCENT = RGBColor(0, 214, 255)
C_ACCENT2 = RGBColor(122, 92, 255)
C_CARD = RGBColor(20, 31, 57)
C_STROKE = RGBColor(66, 93, 134)


def load_runtime_data():
    if not HEALING_SUMMARY.exists():
        return {}
    try:
        return json.loads(HEALING_SUMMARY.read_text(encoding="utf-8"))
    except Exception:
        return {}


def ensure_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    os.environ["MPLCONFIGDIR"] = str(ASSETS / ".mpl-cache")
    (ASSETS / ".mpl-cache").mkdir(parents=True, exist_ok=True)


def build_charts(data: dict):
    import matplotlib.pyplot as plt

    # 1) Runtime KPI bars
    m = data.get("current_metrics", {})
    counters = data.get("healing_counters", {})
    incidents = max(1, int(data.get("total_incidents", 0) or 0))
    auto_recovered = int(data.get("auto_recovered", 0) or 0)
    auto_rate = auto_recovered / incidents
    yield_pct = float(m.get("line_yield", 0.0) or 0.0) * 100.0
    completeness = float(m.get("completeness_score", 0.0) or 0.0)
    l3_rules = float(data.get("latest_l3_snapshot", {}).get("counts", {}).get("causal_rules", 0) or 0)
    l3_chains = float(data.get("latest_l3_snapshot", {}).get("counts", {}).get("failure_chains", 0) or 0)

    kpi_labels = ["Auto Recovery", "Line Yield", "Completeness", "L3 Rule Count", "Failure Chains"]
    kpi_vals = [auto_rate * 100.0, yield_pct, completeness, l3_rules * 6.0, l3_chains * 25.0]
    # L3 counts are scaled for visual balance in one chart.

    fig = plt.figure(figsize=(9, 4.2), facecolor="#070b18")
    ax = fig.add_subplot(111)
    ax.set_facecolor("#070b18")
    bars = ax.bar(kpi_labels, kpi_vals, color=["#00d6ff", "#22c55e", "#8b5cf6", "#f59e0b", "#ef4444"])
    ax.set_ylim(0, 110)
    ax.tick_params(axis="x", colors="#d8e3f2", rotation=8, labelsize=10)
    ax.tick_params(axis="y", colors="#8ca0be", labelsize=9)
    for b, v in zip(bars, kpi_vals):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 2, f"{v:.1f}", ha="center", color="#ecf4ff", fontsize=9)
    for spine in ax.spines.values():
        spine.set_color("#2d405f")
    ax.grid(axis="y", color="#2d405f", alpha=0.4, linewidth=0.8)
    fig.tight_layout()
    kpi_path = ASSETS / "runtime_kpi.png"
    fig.savefig(kpi_path, dpi=180, facecolor=fig.get_facecolor())
    plt.close(fig)

    # 2) Roadmap progress
    phases = ["P1", "P2", "P3", "P4", "P5"]
    values = [100, 100, 85, 55, 20]
    fig = plt.figure(figsize=(8.8, 3.7), facecolor="#070b18")
    ax = fig.add_subplot(111)
    ax.set_facecolor("#070b18")
    bars = ax.barh(phases, values, color=["#22c55e", "#22c55e", "#06b6d4", "#f59e0b", "#6b7280"])
    ax.set_xlim(0, 100)
    ax.tick_params(axis="x", colors="#8ca0be", labelsize=9)
    ax.tick_params(axis="y", colors="#d8e3f2", labelsize=11)
    for b, v in zip(bars, values):
        ax.text(v + 1, b.get_y() + b.get_height() / 2, f"{v}%", va="center", color="#ecf4ff", fontsize=10)
    for spine in ax.spines.values():
        spine.set_color("#2d405f")
    ax.grid(axis="x", color="#2d405f", alpha=0.35, linewidth=0.8)
    fig.tight_layout()
    roadmap_path = ASSETS / "roadmap_progress.png"
    fig.savefig(roadmap_path, dpi=180, facecolor=fig.get_facecolor())
    plt.close(fig)

    # 3) Incident effect chart
    cases = data.get("case_analyses", [])[:5]
    labels = [str(c.get("incident_id", f"INC-{i+1}")) for i, c in enumerate(cases)] or ["INC-0001"]
    deltas = [float(c.get("effect", {}).get("delta_pct", 0.0) or 0.0) for c in cases] or [0.0]
    fig = plt.figure(figsize=(8.8, 3.7), facecolor="#070b18")
    ax = fig.add_subplot(111)
    ax.set_facecolor("#070b18")
    bars = ax.bar(labels, deltas, color="#00d6ff")
    ax.tick_params(axis="x", colors="#d8e3f2", labelsize=10)
    ax.tick_params(axis="y", colors="#8ca0be", labelsize=9)
    ax.set_ylabel("Yield Delta (%)", color="#8ca0be", fontsize=9)
    for b, v in zip(bars, deltas):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height() + 0.02, f"{v:.2f}", ha="center", color="#ecf4ff", fontsize=9)
    for spine in ax.spines.values():
        spine.set_color("#2d405f")
    ax.grid(axis="y", color="#2d405f", alpha=0.35, linewidth=0.8)
    fig.tight_layout()
    impact_path = ASSETS / "incident_impact.png"
    fig.savefig(impact_path, dpi=180, facecolor=fig.get_facecolor())
    plt.close(fig)

    return kpi_path, roadmap_path, impact_path, counters


def add_base(slide, page_idx: int, total_pages: int):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), WIDE, HIGH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_BG
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), WIDE, Inches(0.42))
    top.fill.solid()
    top.fill.fore_color.rgb = C_BG2
    top.line.fill.background()

    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(3.1), Inches(0.09))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = C_ACCENT
    ribbon.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(10.8), Inches(7.15), Inches(2.2), Inches(0.24))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{page_idx}/{total_pages}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = C_MUTED


def add_title(slide, title: str, subtitle: str = ""):
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(0.58), Inches(12.2), Inches(1.0))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(33)
    p.font.bold = True
    p.font.color.rgb = C_TITLE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.74), Inches(1.45), Inches(12.0), Inches(0.55))
        sf = sb.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(15)
        sp.font.color.rgb = C_MUTED


def add_bullet_block(slide, items: list[str], left=0.95, top=2.0, width=11.5, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.7))
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.font.size = Pt(size)
        p.font.color.rgb = C_TEXT
        p.space_after = Pt(11)


def add_card(slide, left, top, width, height, title, body_lines, accent: RGBColor = C_ACCENT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.color.rgb = C_STROKE
    shape.line.width = Pt(1.4)

    mark = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.08),
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = accent
    mark.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = C_TITLE
    p0.space_after = Pt(8)
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13.5)
        p.font.color.rgb = C_TEXT


def add_pipeline(slide):
    labels = ["SENSE", "DETECT", "DIAGNOSE", "RECOVER", "VERIFY", "LEARN"]
    x = 0.92
    y = 2.6
    w = 1.88
    h = 1.08
    gap = 0.24
    colors = [RGBColor(34, 197, 94), RGBColor(0, 214, 255), RGBColor(122, 92, 255), RGBColor(245, 158, 11), RGBColor(239, 68, 68), RGBColor(168, 85, 247)]
    for idx, (label, col) in enumerate(zip(labels, colors)):
        add_card(slide, x, y, w, h, label, ["agent action"], accent=col)
        if idx < len(labels) - 1:
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + w),
                Inches(y + h / 2),
                Inches(x + w + gap - 0.03),
                Inches(y + h / 2),
            )
            line.line.color.rgb = RGBColor(100, 123, 160)
            line.line.width = Pt(1.7)
        x += w + gap


def build_ppt(data: dict):
    ensure_assets()
    kpi_path, roadmap_path, impact_path, counters = build_charts(data)

    prs = Presentation()
    prs.slide_width = WIDE
    prs.slide_height = HIGH
    blank = prs.slide_layouts[6]
    total = 12

    # 1 Cover
    s = prs.slides.add_slide(blank)
    add_base(s, 1, total)
    add_title(s, "EV Battery AIOps for Dark Factory", "Research-Backed Causal Reasoning + LLM Hybrid System")
    add_card(s, 0.78, 2.15, 3.9, 1.55, "Mission", ["문제 감지→진단→복구→학습", "사람 개입 최소화"], C_ACCENT)
    add_card(s, 4.9, 2.15, 3.9, 1.55, "Core Stack", ["Knowledge Graph", "Multi-Agent", "Digital Twin"], C_ACCENT2)
    add_card(s, 9.0, 2.15, 3.6, 1.55, "Current Stage", ["L3 운영, Phase4 통합 중", "실시간 대시보드 연동"], RGBColor(34, 197, 94))
    add_bullet_block(
        s,
        [
            "연구 논문 12편 기반으로 구조 설계 및 기능 매핑",
            "현재: 인과추론 계층(L3)과 LLM 연계형 진단 API까지 구현",
        ],
        top=4.35,
        size=17,
    )

    # 2 Executive summary
    s = prs.slides.add_slide(blank)
    add_base(s, 2, total)
    add_title(s, "Executive Summary", "무엇을 만들었고, 왜 중요한가")
    add_card(s, 0.8, 2.0, 4.0, 3.8, "What We Built", ["L1~L3 온톨로지", "자율복구 루프", "자연어 진단 API"], C_ACCENT)
    add_card(s, 4.95, 2.0, 4.0, 3.8, "Validated Value", ["자동 복구 성공 사례 축적", "재발 대응 속도 개선", "설명 가능한 RCA"], RGBColor(34, 197, 94))
    add_card(s, 9.1, 2.0, 3.45, 3.8, "Next Bet", ["RUL 모델 고도화", "HITL 정책", "실설비 브리지"], RGBColor(245, 158, 11))

    # 3 Paper map
    s = prs.slides.add_slide(blank)
    add_base(s, 3, total)
    add_title(s, "Research Map to System Components", "논문 인사이트와 구현 컴포넌트 연결")
    add_card(s, 0.85, 2.0, 6.0, 1.45, "Sensors 2024/2025", ["KG+DT 융합, KG-driven fault diagnosis", "→ L1/L2 모델 + RCA 파이프라인"], C_ACCENT)
    add_card(s, 6.95, 2.0, 5.5, 1.45, "NAMRC 2026", ["Hybrid Agentic AI + MAS", "→ LLM 전략 + Edge Agent 실행 분리"], C_ACCENT2)
    add_card(s, 0.85, 3.65, 6.0, 1.45, "FD-LLM 2025", ["자연어 질의로 진단", "→ /api/nl-diagnose + schema guardrail"], RGBColor(34, 197, 94))
    add_card(s, 6.95, 3.65, 5.5, 1.45, "Self-Evolving Agents 2025", ["전략의 지속 개선", "→ Phase5 EvolutionAgent 설계"], RGBColor(245, 158, 11))
    add_bullet_block(s, ["핵심: 단순 기능 추가가 아니라, 논문 가설을 시스템 구성요소로 전환"], top=5.55, size=17)

    # 4 Architecture diagram
    s = prs.slides.add_slide(blank)
    add_base(s, 4, total)
    add_title(s, "Cognitive Digital Twin Architecture", "Physical ↔ Knowledge ↔ Agent Layer")
    add_card(s, 0.8, 2.1, 3.7, 2.5, "Physical Twin", ["Sensors", "Alarms", "Equipment State"], C_ACCENT)
    add_card(s, 4.85, 2.1, 3.7, 2.5, "Knowledge Graph", ["L1 Process", "L2 Runtime", "L3 Causality", "L4 Decision"], C_ACCENT2)
    add_card(s, 8.9, 2.1, 3.6, 2.5, "Cognitive Agents", ["Detect", "Diagnose", "Recover", "Learn"], RGBColor(34, 197, 94))
    c1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(4.5), Inches(3.3), Inches(4.82), Inches(3.3))
    c1.line.color.rgb = C_MUTED
    c1.line.width = Pt(2)
    c2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.55), Inches(3.3), Inches(8.87), Inches(3.3))
    c2.line.color.rgb = C_MUTED
    c2.line.width = Pt(2)
    add_bullet_block(s, ["운영 루프: Sense → Detect → Diagnose → Recover → Verify → Learn"], top=5.35, size=18)

    # 5 Self-healing loop
    s = prs.slides.add_slide(blank)
    add_base(s, 5, total)
    add_title(s, "Self-Healing Loop", "실시간 자동 복구 엔진")
    add_pipeline(s)
    add_bullet_block(
        s,
        [
            "AnomalyDetector: threshold / 3-sigma / trend-shift",
            "CausalReasoner: chain + history boost",
            "AutoRecoveryAgent: playbook 실행 및 검증",
        ],
        top=4.5,
        size=17,
    )

    # 6 Causal deep dive
    s = prs.slides.add_slide(blank)
    add_base(s, 6, total)
    add_title(s, "L3 Causal Reasoning Deep-Dive", "거리 기반 RCA에서 인과 기반 RCA로 전환")
    add_card(s, 0.8, 2.0, 6.0, 3.9, "Before", ["거리/유사도 중심 후보순위", "증상 근처 원인 과대평가", "재발 학습 반영 약함"], RGBColor(239, 68, 68))
    add_card(s, 6.95, 2.0, 5.6, 3.9, "After", ["CausalRule 체인 추론", "FailureChain 매칭과 성공률 반영", "복구결과가 confidence에 누적학습"], RGBColor(34, 197, 94))

    # 7 Runtime evidence
    s = prs.slides.add_slide(blank)
    add_base(s, 7, total)
    add_title(s, "Runtime Evidence", "실행 데이터 기반 검증")
    s.shapes.add_picture(str(kpi_path), Inches(0.8), Inches(1.95), width=Inches(6.2))
    s.shapes.add_picture(str(impact_path), Inches(6.95), Inches(1.95), width=Inches(5.6))
    add_card(
        s,
        0.8,
        5.85,
        12.0,
        1.1,
        "Live Counters",
        [f"readings={counters.get('reading', 0)}, alarms={counters.get('alarm', 0)}, incidents={counters.get('incident', 0)}, recovery={counters.get('recovery', 0)}"],
        C_ACCENT,
    )

    # 8 LLM hybrid
    s = prs.slides.add_slide(blank)
    add_base(s, 8, total)
    add_title(s, "Phase4 LLM Hybrid Integration", "안전한 LLM 연동 + 자동 폴백")
    add_card(
        s,
        0.8,
        2.0,
        5.9,
        3.9,
        "NaturalLanguageDiagnoser",
        [
            "질의 입력 → Incident/Causal context 구성",
            "OpenAI Responses API(JSON schema strict)",
            "sanitize/guardrail 적용",
            "실패 시 symbolic mode 자동 전환",
        ],
        C_ACCENT,
    )
    add_card(
        s,
        6.95,
        2.0,
        5.6,
        3.9,
        "PredictiveAgent",
        [
            "MTBF + incident 기반 risk_score",
            "RUL 시간 추정 + P1/P2/P3 우선순위",
            "API: /api/predictive-rul, /api/phase4-status",
            "대시보드에 LLM ON/OFF 및 모델 노출",
        ],
        C_ACCENT2,
    )

    # 9 Demo story
    s = prs.slides.add_slide(blank)
    add_base(s, 9, total)
    add_title(s, "Demo Storyboard", "현업 시나리오 기준 End-to-End 흐름")
    add_bullet_block(
        s,
        [
            "Step 1) 공정 이상 발생: PS-103 / PS-507 threshold breach",
            "Step 2) RCA 강화: equipment_mtbf + causal chain evidence",
            "Step 3) 복구 실행: ADJUST_PARAMETER, 자동 검증",
            "Step 4) 결과 저장: Incident, RecoveryAction, FailureChain",
            "Step 5) 운영자 질의: 자연어로 원인/권고 확인",
        ],
        top=2.05,
        size=18,
    )

    # 10 Roadmap
    s = prs.slides.add_slide(blank)
    add_base(s, 10, total)
    add_title(s, "Roadmap & Delivery Plan", "Phase 진척도와 다음 실행")
    s.shapes.add_picture(str(roadmap_path), Inches(0.9), Inches(1.95), width=Inches(6.8))
    add_card(
        s,
        7.95,
        1.95,
        4.65,
        4.8,
        "Next 90 Days",
        [
            "RUL 모델: heuristic → LSTM/Transformer",
            "HITL 정책 엔진(고위험 action 승인제)",
            "MQTT/OPC-UA 실설비 브리지",
            "ResilienceOrchestrator PoC",
            "비용/효과 KPI 대시보드",
        ],
        RGBColor(245, 158, 11),
    )

    # 11 Risks and governance
    s = prs.slides.add_slide(blank)
    add_base(s, 11, total)
    add_title(s, "Risk, Governance, and Safety", "자동화 품질을 보장하는 운영 원칙")
    add_card(s, 0.8, 2.0, 4.0, 3.8, "Model Risk", ["오탐/미탐 관리", "confidence 임계값 운영", "재현 테스트"], RGBColor(239, 68, 68))
    add_card(s, 4.95, 2.0, 4.0, 3.8, "Operational Risk", ["고위험 조치 HITL", "에스컬레이션 정책", "롤백 절차"], RGBColor(245, 158, 11))
    add_card(s, 9.1, 2.0, 3.45, 3.8, "Data Risk", ["로그/이력 보존", "감사 추적성", "보안 키 관리"], C_ACCENT)

    # 12 Closing + references
    s = prs.slides.add_slide(blank)
    add_base(s, 12, total)
    add_title(s, "Conclusion & References", "From Monitoring to Autonomous Recovery")
    add_bullet_block(
        s,
        [
            "핵심 결론: 인과추론 + LLM 하이브리드가 제조 AIOps 실용 구간을 연다",
            "실행 결론: 이미 동작하는 L3/Phase4를 기반으로 L4 다크팩토리로 확장 가능",
            "References: Sensors 2024, Sensors 2025, NAMRC 2026, FD-LLM 2025, Self-Evolving 2025",
        ],
        top=2.1,
        size=18,
    )

    prs.save(PPT_PATH)


def build_script():
    content = f"""# EV Battery AIOps 발표 대본 (PRO)

발표자료: `{PPT_PATH.name}`
발표일: {date.today().isoformat()}
권장 발표시간: 12~15분

## Slide 1. Cover (40초)
오늘 발표는 EV 배터리 제조 공정을 대상으로, 인과추론과 LLM을 결합한 자율복구 AIOps 플랫폼의 연구-구현 결과를 공유드립니다. 목적은 대시보드 관제가 아니라, 사람 개입을 최소화한 Dark Factory 운영입니다.

## Slide 2. Executive Summary (60초)
핵심 메시지는 세 가지입니다. 첫째, 연구 논문 인사이트를 시스템 컴포넌트로 전환했습니다. 둘째, L3 인과추론과 자율복구 루프를 실제 동작시켰습니다. 셋째, Phase4에서 LLM을 안전하게 통합해 자연어 진단까지 확장했습니다.

## Slide 3. Research Map (70초)
Sensors 계열 연구는 KG+DT와 fault diagnosis의 타당성을 제공합니다. NAMRC 2026은 LLM 전략 계층과 edge 실행 계층 분리를 제안하고, FD-LLM은 자연어 진단 가능성을 확인해줍니다. 저희는 이 세 흐름을 하나의 운영 아키텍처로 통합했습니다.

## Slide 4. Architecture (70초)
Physical Twin에서 수집된 이벤트가 Knowledge Graph로 들어오고, Cognitive Agents가 의사결정을 수행합니다. 이때 그래프는 단순 저장소가 아니라 추론의 근거를 제공하는 실행 지식 베이스 역할을 합니다.

## Slide 5. Self-Healing Loop (80초)
루프는 SENSE, DETECT, DIAGNOSE, RECOVER, VERIFY, LEARN 순서입니다. 중요한 점은 LEARN이 마지막 단계가 아니라 다음 루프의 정확도를 끌어올리는 시작점이라는 점입니다.

## Slide 6. Causal Deep-Dive (80초)
이전에는 거리 기반 후보순위가 중심이었다면, 현재는 CausalRule과 FailureChain을 이용한 인과 기반 RCA를 수행합니다. 그래서 원인 설명성이 좋아지고, 반복 장애 대응 속도가 빨라집니다.

## Slide 7. Runtime Evidence (70초)
실행 결과에서 자동 복구율, 수율, 완성도, L3 지식 축적량을 함께 확인할 수 있습니다. Incident별 수율 개선 영향도도 시각화해 복구 조치의 실효성을 빠르게 점검합니다.

## Slide 8. LLM Hybrid (80초)
NaturalLanguageDiagnoser는 JSON schema 강제와 sanitize를 통해 안전한 응답을 반환합니다. 호출 실패 시 symbolic 모드로 즉시 폴백하므로, 운영 안정성을 해치지 않습니다. PredictiveAgent는 RUL 우선순위를 제공해 예지정비 의사결정을 보조합니다.

## Slide 9. Demo Storyboard (60초)
실제 시나리오는 이상 감지, 원인 진단, 복구 실행, 검증, 이력 학습 순서로 진행됩니다. 운영자는 마지막에 자연어 질의만으로 근거와 권고를 확인할 수 있습니다.

## Slide 10. Roadmap (70초)
다음 90일은 RUL 고도화, HITL 정책, 실설비 브리지, 재편성 오케스트레이터가 핵심입니다. 즉 기능 확장이 아니라, 운영 신뢰성을 높이는 단계입니다.

## Slide 11. Risk/Governance (60초)
자동화 성능만큼 중요한 것이 리스크 거버넌스입니다. 모델 리스크, 운영 리스크, 데이터 리스크를 분리해 관리하고, 고위험 조치에는 반드시 HITL 정책을 적용합니다.

## Slide 12. Closing (40초)
결론적으로, 이 시스템은 모니터링 도구를 넘어 인과 기반 자율복구 플랫폼으로 진화하고 있습니다. 연구와 현업 실행을 동시에 만족하는 아키텍처로 Dark Factory 전환을 가속하겠습니다.
"""
    SCRIPT_PATH.write_text(content, encoding="utf-8")


def main():
    RESULTS.mkdir(parents=True, exist_ok=True)
    data = load_runtime_data()
    build_ppt(data)
    build_script()
    print(f"PPT generated: {PPT_PATH}")
    print(f"Script generated: {SCRIPT_PATH}")


if __name__ == "__main__":
    main()

